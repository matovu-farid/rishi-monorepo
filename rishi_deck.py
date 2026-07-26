from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE

OUT = "/private/tmp/Rishi-Google-AI-Lab-Pitch-Deck-First-Draft.pptx"
LOGO = "/Users/faridmatovu/projects/rishi-monorepo/apps/apple/rishi/rishi/Assets.xcassets/rishi.imageset/rishi-icon (2).png"
IPHONE_LIBRARY = "/Users/faridmatovu/Desktop/rishi screenshots/Simulator Screenshot - iPhone 12 Pro Max - 2026-07-25 at 12.19.40.png"
IPHONE_READER = "/Users/faridmatovu/Desktop/rishi screenshots/Simulator Screenshot - iPhone 12 Pro Max - 2026-07-25 at 12.17.22.png"
IPHONE_TTS = "/Users/faridmatovu/Desktop/rishi screenshots/Simulator Screenshot - iPhone 12 Pro Max - 2026-07-25 at 12.18.49.png"
DEMO_POSTER = "/Users/faridmatovu/projects/rishi-monorepo/deck-assets/rishi-demo-poster.png"

W, H = 13.333, 7.5
NAVY = RGBColor(14, 28, 41)
NAVY2 = RGBColor(22, 43, 61)
PAPER = RGBColor(246, 240, 224)
BLUE = RGBColor(98, 158, 199)
MUTED = RGBColor(164, 185, 198)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

def fill_bg(slide):
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = NAVY

def add_text(slide, text, x, y, w, h, size=20, color=PAPER, bold=False, font="Aptos", align=PP_ALIGN.LEFT, margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    return box

def add_logo(slide, x=12.28, y=0.3, size=0.55):
    slide.shapes.add_picture(LOGO, Inches(x), Inches(y), width=Inches(size), height=Inches(size))

def add_footer(slide, num):
    add_text(slide, "RISHI  •  GOOGLE AI LAB  •  FIRST DRAFT", 0.68, 7.08, 4.5, 0.2, 8, BLUE, True)
    add_text(slide, f"{num:02d}", 12.2, 7.05, 0.45, 0.2, 9, MUTED, True, align=PP_ALIGN.RIGHT)

def add_header(slide, kicker, title, num):
    add_logo(slide)
    add_text(slide, kicker.upper(), 0.7, 0.42, 5.5, 0.25, 10, BLUE, True)
    add_text(slide, title, 0.7, 0.78, 9.8, 0.75, 27, PAPER, True, "Aptos Display")
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.72), Inches(0.36), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    add_footer(slide, num)

def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = NAVY2; card.line.color.rgb = RGBColor(39, 68, 88)
    add_text(slide, title, x+0.25, y+0.22, w-0.5, 0.3, 14, accent, True)
    add_text(slide, body, x+0.25, y+0.68, w-0.5, h-0.85, 16, PAPER)

def add_bullets(slide, items, x, y, w, h, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0; p.space_after = Pt(11)
        p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = PAPER
        p._p.get_or_add_pPr().insert(0, p._p._new_buChar()) if False else None
    return box

# 1 Cover
s = prs.slides.add_slide(blank); fill_bg(s); add_logo(s, 0.75, 0.65, 0.82)
add_text(s, "RISHI", 0.75, 1.7, 4, 0.4, 15, BLUE, True)
add_text(s, "A teacher\ninside every book.", 0.72, 2.12, 7.7, 1.55, 38, PAPER, True, "Aptos Display")
add_text(s, "Make dense material easier to understand through explanations, conversation, reading, and listening.", 0.78, 4.15, 6.4, 0.8, 19, MUTED)
add_text(s, "Google AI Lab application  •  First draft", 0.78, 6.55, 5, 0.3, 11, BLUE, True)
add_text(s, "01", 12.2, 7.05, 0.45, 0.2, 9, MUTED, True, align=PP_ALIGN.RIGHT)

# 2 Problem
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "The problem", "A book often assumes the reader already has a teacher.", 2)
add_text(s, "Dense material can be difficult to follow, break down, remember, and discuss.", 0.75, 2.1, 6.1, 0.85, 24, PAPER, True)
add_text(s, "When a learner gets stuck, progress often stops—not because they lack motivation, but because the explanation is missing.", 0.78, 3.3, 5.9, 0.95, 18, MUTED)
add_card(s, 7.55, 2.05, 4.75, 2.0, "THE QUESTION", "What if the guidance you need lived inside the book?", BLUE)
add_text(s, "Learning should not stop at the moment a reader gets stuck.", 0.78, 4.9, 7.6, 0.45, 18, MUTED)

# 3 Idea + iPhone screen
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "The idea", "A teacher-like assistant, inside the reading experience.", 3)
add_text(s, "Ask for an explanation. Break down a difficult passage. Explore the meaning in context.", 0.75, 2.1, 5.0, 1.0, 21, PAPER, True)
add_text(s, "Rishi turns the moment of confusion into a next step.", 0.75, 3.45, 5.0, 0.55, 20, BLUE, True)
add_text(s, "The goal is simple: make learning easier for anyone with a book.", 0.75, 4.05, 4.9, 0.75, 17, MUTED)
s.shapes.add_picture(IPHONE_LIBRARY, Inches(8.0), Inches(1.95), height=Inches(4.75))

# 4 Two ways
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "Two ways to learn", "See the material. Hear the material. Keep moving.", 4)
add_card(s, 0.8, 2.1, 5.55, 2.5, "READ", "Follow the page, highlight ideas, search the text, and stay in context.", BLUE)
add_card(s, 6.95, 2.1, 5.55, 2.5, "LISTEN", "Use text-to-speech to hear the material while reading, walking, commuting, or reviewing.", RGBColor(185, 147, 88))
add_text(s, "Rishi makes learning more flexible by bringing visual and auditory modes together.", 0.8, 5.4, 8.5, 0.45, 18, MUTED)
s.shapes.add_picture(IPHONE_TTS, Inches(10.25), Inches(0.55), height=Inches(5.8))

# 5 AI
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "Contextual AI", "The assistant is grounded in the book you are reading.", 5)
add_text(s, "Rishi is designed to help a reader move from “I’m stuck” to “I understand” without leaving the material.", 0.75, 2.1, 5.2, 1.2, 21, PAPER, True)
add_text(s, "This is not about replacing teachers. It is about making guidance available in the moment a learner needs it.", 0.75, 4.25, 5.4, 0.85, 17, MUTED)
add_bullets(s, ["Book-specific conversations", "Text and semantic/vector search", "Retrieval and embeddings/indexing", "Explanations connected to current material"], 7.0, 2.1, 5.1, 2.8, 17)

# 6 Product demo
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "Product demo", "A 68-second journey from page to assisted learning.", 6)
add_text(s, "The recording shows the real iPhone experience: open a book, stay with the page, listen, and keep moving.", 0.75, 2.05, 5.3, 1.0, 20, PAPER, True)
add_card(s, 0.8, 3.65, 5.0, 1.35, "WATCH THE RECORDING", "[Insert public demo-video link]", RGBColor(185, 147, 88))
s.shapes.add_picture(DEMO_POSTER, Inches(8.0), Inches(1.95), height=Inches(4.75))

# 7 Workflow + iPhone screen
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "From book to understanding", "Turn reading into an active learning workflow.", 7)
add_text(s, "Choose a book  →  Read and listen  →  Ask for an explanation\nSearch and revisit  →  Highlight and bookmark the insight", 0.75, 2.1, 5.8, 1.15, 19, PAPER, True)
add_text(s, "The product brings the teacher, the page, and the voice into one continuous experience.", 0.75, 4.35, 5.5, 0.75, 16, MUTED)
s.shapes.add_picture(IPHONE_READER, Inches(8.0), Inches(1.95), height=Inches(4.75))

# 8 Apple-first wedge
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "Go-to-market", "Win Apple first. Expand from a position of strength.", 8)
add_text(s, "Rishi’s first ecosystem is Apple: iPhone for everyday learning, iPad for deeper reading, and App Store distribution for a focused launch.", 0.75, 2.0, 6.1, 1.15, 22, PAPER, True)
add_card(s, 0.8, 3.65, 3.55, 1.55, "NOW", "iPhone\nApp Store launch")
add_card(s, 4.85, 3.65, 3.55, 1.55, "NEXT", "iPad\nLong-form reading")
add_card(s, 8.9, 3.65, 3.55, 1.55, "LATER", "Other platforms\nAfter Apple fit")
add_text(s, "[Add exact App Store status and date]", 0.8, 5.75, 4.5, 0.3, 12, MUTED)

# 9 Metrics
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "Measurement", "We are pre-revenue. Launch is where the learning begins.", 9)
add_text(s, "We do not yet have meaningful usage or validation metrics. At launch, we will establish baselines for:", 0.75, 2.05, 7.0, 0.9, 20, PAPER, True)
add_card(s, 0.8, 3.35, 2.85, 1.6, "01", "Activation\nand books imported")
add_card(s, 3.9, 3.35, 2.85, 1.6, "02", "Reading and\nlistening sessions")
add_card(s, 7.0, 3.35, 2.85, 1.6, "03", "AI questions\nand explanations")
add_card(s, 10.1, 3.35, 2.3, 1.6, "04", "Retention and\nfeedback")

# 10 Why now
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "Why Rishi / why now", "Knowledge should feel less solitary—and Apple is the right first home.", 10)
add_text(s, "Conversational AI can make the act of learning feel guided, responsive, and human.", 0.75, 2.15, 6.2, 0.8, 24, PAPER, True)
add_text(s, "Rishi brings that guidance directly into the book—so the learner can ask, listen, revisit, and understand in context.", 0.78, 3.75, 6.0, 1.0, 19, MUTED)
add_card(s, 7.55, 2.05, 4.65, 1.7, "APPLE-FIRST ADVANTAGE", "A coherent iPhone + iPad reading experience, shipped and refined before expanding.", BLUE)
add_card(s, 7.55, 4.1, 4.65, 1.35, "FOUNDER STORY", "[Add founder story and location]", RGBColor(185, 147, 88))

# 11 Ask
s = prs.slides.add_slide(blank); fill_bg(s); add_header(s, "The ask", "Help us make understanding easier for anyone with a book.", 11)
add_text(s, "We are seeking support to turn a launch-ready prototype into a reliable daily learning tool.", 0.75, 2.05, 7.0, 0.75, 23, PAPER, True)
add_bullets(s, ["AI product and model guidance", "Evaluation expertise for explanations and grounded answers", "Launch mentorship and feedback"], 0.85, 3.45, 6.4, 2.1, 18)
add_card(s, 8.1, 2.65, 4.2, 1.95, "NEXT STEP", "[Refine the specific Google AI Lab request]", RGBColor(185, 147, 88))

prs.save(OUT)
print(OUT)
